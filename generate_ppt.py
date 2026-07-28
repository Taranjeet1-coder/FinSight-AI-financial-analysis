import collections 
import collections.abc
# python-pptx needs collections.abc which changed in 3.10
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Colors
    NAVY = RGBColor(20, 32, 72)
    GOLD = RGBColor(201, 162, 39)
    WHITE = RGBColor(255, 255, 255)
    
    # Layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY
        
    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(title_slide_layout)
        set_slide_background(slide)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = WHITE
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = subtitle_text
        for p in subtitle.text_frame.paragraphs:
            p.font.color.rgb = GOLD
            p.font.size = Pt(24)
            
        return slide

    def add_bullet_slide(title_text, bullets):
        slide = prs.slides.add_slide(bullet_slide_layout)
        set_slide_background(slide)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title_text
        title_shape.text_frame.paragraphs[0].font.color.rgb = GOLD
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        tf = body_shape.text_frame
        tf.clear()
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.color.rgb = WHITE
            p.font.size = Pt(20)
            p.level = 0
            
        return slide

    # Slide 1
    add_title_slide(
        "Financial Intelligence System for Comparative Analysis of Bitcoin, Gold, and the S&P 500 using Artificial Intelligence",
        "An Advanced Exploratory Data Analysis & Machine Learning Approach\n\nTeam Members: [Names]\nGuide: [Guide Name]\nUniversity: [University Name]"
    )

    # Slide 2
    add_bullet_slide("Executive Summary", [
        "Analyzed 11 years of daily financial data to compare crypto versus traditional asset performance.",
        "Applied classical Machine Learning, Deep Learning (LSTM/GRU), and Soft Computing techniques.",
        "Discovered that classical ML and Fuzzy Logic provide robust signals compared to deep learning noise.",
        "Developed a fully interactive Streamlit dashboard for real-time portfolio optimization and regime detection."
    ])

    # Slide 3
    add_bullet_slide("Project Objectives & Scope", [
        "Compare Asset Classes: Evaluate Bitcoin's role against traditional safe havens (Gold) and equities (S&P 500).",
        "Quantify Risk: Model volatility using advanced econometric frameworks (GARCH, VaR).",
        "Predictive Modeling: Assess the viability of LSTM/GRU networks versus Random Forest/XGBoost.",
        "Optimization: Apply Genetic Algorithms and Markowitz theory to find optimal portfolio weights."
    ])

    # Slide 4
    add_bullet_slide("Analytical Methodology Framework", [
        "Phase 1 - Data Engineering: Cleaning, scaling, and handling missing values.",
        "Phase 2 - Exploratory Analysis & Risk: Descriptive stats, correlation, and volatility clustering.",
        "Phase 3 - Advanced Modeling: Machine Learning, Deep Learning, and Soft Computing.",
        "Phase 4 - Optimization & Deployment: Portfolio construction and Streamlit application."
    ])

    # Slide 5
    add_bullet_slide("Data Overview & Data Cleaning", [
        "11 years of daily frequency data (mid-2014 to 2025), yielding ~2,800 valid observations.",
        "Addressed missing values using forward-fill and backward-fill techniques to maintain integrity.",
        "Computed logarithmic returns and 30-day/60-day rolling statistics.",
        "Verified stationarity across all return series using Augmented Dickey-Fuller (ADF) tests."
    ])

    # Slide 6
    add_bullet_slide("Descriptive Statistics", [
        "Bitcoin exhibits extreme kurtosis (~13), indicating a high frequency of 'fat tail' events.",
        "S&P 500 shows consistent steady growth with moderate skewness.",
        "Gold acts as a stable anchor with lower volatility and standard deviation.",
        "Non-normal return distributions observed across all assets, requiring robust risk metrics."
    ])

    # Slide 7
    add_bullet_slide("Visual Trends: Normalized Price Evolution", [
        "Significant outperformance of Bitcoin in cumulative growth, but with severe drawdowns.",
        "Steady, compounding upward trajectory for the S&P 500.",
        "Gold's performance remains independent of equity and crypto market exuberance.",
        "Identifiable regime shifts during macroeconomic shocks (e.g., COVID-19)."
    ])

    # Slide 8
    add_bullet_slide("Correlation & Macro Relationships", [
        "Minimal to negative correlation between Gold and equities, reinforcing its hedge status.",
        "Bitcoin's correlation with the S&P 500 has varied, challenging its 'digital gold' narrative.",
        "Bitcoin averaged 0.49% returns in high-inflation periods vs. 0.11% in low-inflation.",
        "Rolling 60-day correlations reveal dynamic relationships dependent on market regimes."
    ])

    # Slide 9
    add_bullet_slide("Risk Analysis: Asset Drawdowns", [
        "Bitcoin experienced multiple drawdowns exceeding 70%, reflecting extreme risk.",
        "S&P 500 drawdowns were contained within expected bounds, barring systemic shocks.",
        "Value-at-Risk (VaR) and Conditional VaR (CVaR) metrics highlight significant tail risk in crypto.",
        "[Placeholder for drawdowns.png]"
    ])

    # Slide 10
    add_bullet_slide("Volatility Modeling: GARCH Framework", [
        "Applied GARCH(1,1) modeling to capture volatility clustering.",
        "Results show strong volatility persistence (alpha + beta ≈ 0.92).",
        "Market shocks have a long-lasting effect on future volatility.",
        "Crucial for adjusting portfolio weights dynamically during turbulent periods."
    ])

    # Slide 11
    add_bullet_slide("Machine Learning: Regime Classification", [
        "Classified market regimes based on volatility and returns using Random Forest.",
        "Achieved 89.49% classification accuracy.",
        "Successfully identified distinct periods of bull, bear, and transitional markets.",
        "VIX levels and rolling volatility emerged as the most critical features."
    ])

    # Slide 12
    add_bullet_slide("Predictive Modeling: Classical ML", [
        "Benchmarked Random Forest and XGBoost against baseline models.",
        "Both achieved modest positive predictive returns (approx. 0.21 - 0.23).",
        "Tree-based algorithms handled non-linear, noisy daily financial returns well.",
        "Simpler, interpretable models can be highly effective in finance."
    ])

    # Slide 13
    add_bullet_slide("Deep Learning Limitations", [
        "Implemented LSTM and GRU architectures for time-series forecasting.",
        "Deep learning models underperformed classical ML and simple mean predictions.",
        "Struggled with market noise, leading to negative R² scores and diverging loss.",
        "Highlights the risk of overfitting neural networks on volatile financial data."
    ])

    # Slide 14
    add_bullet_slide("Soft Computing: Fuzzy Logic Risk Scoring", [
        "Designed a Fuzzy Logic inference system to generate dynamic risk scores.",
        "Generated 2,326 'HOLD' and 486 'SELL' signals (avoided false 'BUY' signals).",
        "Provides a rules-based, interpretable layer to complement black-box ML models.",
        "[Placeholder for fuzzy_risk.png]"
    ])

    # Slide 15
    add_bullet_slide("Soft Computing: Genetic Algorithms", [
        "Deployed a Genetic Algorithm (GA) to discover optimal portfolio weights.",
        "Optimized for maximum Sharpe Ratio across thousands of generations.",
        "Successfully navigated the complex, non-convex optimization landscape.",
        "Validates traditional Mean-Variance optimization results."
    ])

    # Slide 16
    add_bullet_slide("Portfolio Optimization & Efficient Frontier", [
        "Ran Monte Carlo simulations on over 5,000 portfolios to map the Efficient Frontier.",
        "Evaluated Minimum Variance, Risk Parity, and Maximum Sharpe strategies.",
        "Identified the optimal risk-reward tradeoff point for this asset mix.",
        "Proves the necessity of diversification in dampening Bitcoin's volatility."
    ])

    # Slide 17
    add_bullet_slide("Final Optimized Portfolio Allocation", [
        "Gold (56.3%): Acts as the primary stabilizer and risk mitigator.",
        "S&P 500 (38.7%): Provides consistent, compounding equity growth.",
        "Bitcoin (5.0%): Serves as a high-growth satellite asset to boost yield.",
        "Allocation balances extreme upside potential with severe drawdown protection."
    ])

    # Slide 18
    add_bullet_slide("Strategy Backtesting: COVID-19", [
        "Evaluated portfolio performance during the Q1 2020 COVID-19 market crash.",
        "Gold surged +21%, significantly offsetting losses.",
        "S&P 500 dropped -6%, and Bitcoin fell -3% during the measured window.",
        "Diversified strategy outperformed holding pure equities or pure crypto."
    ])

    # Slide 19
    add_bullet_slide("Interactive Dashboard Deployment", [
        "Built a Streamlit dashboard (app.py) for real-time portfolio management.",
        "Live data fetching via yfinance for BTC, SP500, Gold, and VIX.",
        "Core Modules: Market snapshot, 3D Portfolio Efficient Frontier, SIP calculator.",
        "AI Insights: Rule-based ML detection (KMeans) linking VIX to market regimes."
    ])

    # Slide 20
    add_bullet_slide("Key Findings & Business Impact", [
        "Asset Synergies: 5% crypto allocation maximizes Sharpe ratio safely.",
        "Model Reality: Classical ML (Random Forest) is more reliable than Deep Learning here.",
        "Risk Management: GARCH and Fuzzy Logic provide actionable downside protection.",
        "The Streamlit dashboard bridges complex analysis and retail accessibility."
    ])

    # Slide 21
    add_bullet_slide("Conclusion & Future Scope", [
        "Algorithmic portfolio optimization successfully blends traditional and digital assets.",
        "Future Scope: Integrate sentiment analysis from financial news and Twitter/X.",
        "Future Scope: Expand dataset to include bonds, real estate (REITs), and altcoins.",
        "Future Scope: Migrate the dashboard to a cloud-native microservices architecture."
    ])

    # Slide 22
    add_title_slide("Questions & Discussion", "Thank you for your time and attention.\nOpen for Q&A.")

    prs.save("Project_Analysis_Presentation.pptx")
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
